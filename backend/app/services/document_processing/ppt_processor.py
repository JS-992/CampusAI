from pptx import Presentation


def extract_text_from_ppt(file_path: str) -> str:
    """
    Extract text from all slides in a PowerPoint file.

    Args:
        file_path: Path to the PowerPoint file.

    Returns:
        All extracted slide text combined into one string.
    """

    presentation = Presentation(file_path)

    extracted_text = []

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                if shape.text.strip():
                    extracted_text.append(shape.text)

    return "\n".join(extracted_text)